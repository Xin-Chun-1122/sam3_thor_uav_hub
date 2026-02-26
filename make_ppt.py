#!/usr/bin/env python3
"""
make_ppt.py
===========
產生 UAV 雙目視覺事件式告警系統 PPT 報告
執行：python3 make_ppt.py
輸出：UAV_Alert_System_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── 色票 ───
C_BG       = RGBColor(0x1a, 0x1a, 0x2e)   # 深藍背景
C_CARD     = RGBColor(0x16, 0x21, 0x3e)   # 卡片藍
C_ACCENT   = RGBColor(0xe9, 0x45, 0x60)   # 紅色強調
C_ACCENT2  = RGBColor(0x5d, 0xad, 0xe2)   # 淺藍
C_WHITE    = RGBColor(0xff, 0xff, 0xff)
C_GRAY     = RGBColor(0xaa, 0xaa, 0xaa)
C_GREEN    = RGBColor(0x27, 0xae, 0x60)
C_ORANGE   = RGBColor(0xe6, 0x7e, 0x22)
C_YELLOW   = RGBColor(0xf1, 0xc4, 0x0f)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # 完全空白


def add_slide():
    s = prs.slides.add_slide(blank_layout)
    # 深藍背景
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = C_BG
    return s


def box(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=14, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, space_before=0):
    from pptx.util import Pt as Ptx
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Ptx(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Ptx(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def header_bar(slide, title_text, subtitle_text="", slide_num=""):
    """頂部標題列"""
    box(slide, 0, 0, 13.333, 1.1, fill=C_CARD)
    box(slide, 0, 0, 0.25, 1.1, fill=C_ACCENT)   # 左側紅色邊條
    txt(slide, title_text, 0.35, 0.08, 11, 0.55,
        size=32, bold=True, color=C_WHITE)
    if subtitle_text:
        txt(slide, subtitle_text, 0.35, 0.62, 10, 0.42,
            size=16, color=C_ACCENT2)
    if slide_num:
        txt(slide, slide_num, 12.5, 0.35, 0.7, 0.4,
            size=14, color=C_GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, text="無人機雙目視覺事件式告警回傳系統 · NTHU · 2026"):
    box(slide, 0, 7.15, 13.333, 0.35, fill=C_CARD)
    txt(slide, text, 0.3, 7.17, 12, 0.28,
        size=11, color=C_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# Slide 1 — 封面
# ═══════════════════════════════════════════════
s = add_slide()
# 頂部裝飾條
box(s, 0, 0, 13.333, 0.12, fill=C_ACCENT)
# 中央卡片
box(s, 1.0, 1.0, 11.333, 5.5, fill=C_CARD,
    line=C_ACCENT2, line_w=Pt(1.5))
box(s, 1.0, 1.0, 0.3, 5.5, fill=C_ACCENT)

txt(s, "無人機雙目視覺", 1.5, 1.3, 10.5, 1.0,
    size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
txt(s, "事件式告警回傳系統", 1.5, 2.15, 10.5, 1.0,
    size=44, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

box(s, 3.5, 3.3, 6.333, 0.04, fill=C_ACCENT2)

txt(s, "UAV Dual-Camera Event-Driven Alert Transmission System", 1.5, 3.45, 10.5, 0.5,
    size=16, color=C_ACCENT2, align=PP_ALIGN.CENTER)

txt(s, "基於 Meta SAM3 零樣本偵測 + 事件狀態機 + WireGuard VPN", 1.5, 4.0, 10.5, 0.45,
    size=15, color=C_GRAY, align=PP_ALIGN.CENTER)

txt(s, "平台：NVIDIA Jetson AGX Thor  ·  Intel RealSense D435 × 2  ·  ArduPilot MAVLink",
    1.5, 4.5, 10.5, 0.4, size=13, color=C_GRAY, align=PP_ALIGN.CENTER)

box(s, 3.5, 5.05, 6.333, 0.04, fill=C_ACCENT2)

txt(s, "國立清華大學  ·  2026-02-25", 1.5, 5.15, 10.5, 0.4,
    size=13, color=C_GRAY, align=PP_ALIGN.CENTER)

box(s, 0, 7.38, 13.333, 0.12, fill=C_ACCENT)
footer(s)


# ═══════════════════════════════════════════════
# Slide 2 — 研究動機與問題定義
# ═══════════════════════════════════════════════
s = add_slide()
header_bar(s, "研究動機與問題定義",
           "為什麼需要事件式傳輸？", "02")
footer(s)

# 左欄：問題
box(s, 0.3, 1.25, 6.1, 5.6, fill=C_CARD, line=C_ORANGE, line_w=Pt(1))
box(s, 0.3, 1.25, 6.1, 0.42, fill=C_ORANGE)
txt(s, "  ⚠️  傳統連續串流的問題", 0.3, 1.25, 6.1, 0.42,
    size=14, bold=True, color=C_WHITE)

problems = [
    ("頻寬限制", "山區 4G LTE 上行僅 200 Kbps–2 Mbps\n720p@15fps H.264 需 1,000–2,000 Kbps → 不可行"),
    ("假陽性誤報", "雲、光影、枝葉晃動造成間歇性誤偵測\n單幀偵測命中即觸發 → 大量無意義傳輸"),
    ("電力消耗", "持續串流需持續高功率運算\n無人機續航時間受限"),
    ("儲存壓力", "連續錄影產生大量無用影像\n難以即時分析有效事件"),
]
y_p = 1.8
for title, desc in problems:
    box(s, 0.45, y_p, 5.8, 0.8, fill=RGBColor(0x1e, 0x2a, 0x4a), line=C_ORANGE, line_w=Pt(0.5))
    txt(s, f"◆  {title}", 0.55, y_p + 0.03, 5.6, 0.28, size=13, bold=True, color=C_ORANGE)
    txt(s, desc, 0.55, y_p + 0.30, 5.6, 0.45, size=11, color=C_GRAY)
    y_p += 0.98

# 右欄：解決方案
box(s, 6.9, 1.25, 6.1, 5.6, fill=C_CARD, line=C_GREEN, line_w=Pt(1))
box(s, 6.9, 1.25, 6.1, 0.42, fill=C_GREEN)
txt(s, "  ✅  本系統解決方案", 6.9, 1.25, 6.1, 0.42,
    size=14, bold=True, color=C_WHITE)

solutions = [
    ("事件式傳輸", "只在偵測確認後才傳輸\n頻寬節省 95–99%，靜默時 = 0 Kbps"),
    ("狀態機過濾", "5 狀態 FSM 累積命中率\n首次觸發需 ~3 秒，消除瞬間假陽性"),
    ("SAM3 零樣本偵測", "Meta SAM3 text-grounding\n無需重新訓練即可偵測 fire/smoke/person"),
    ("WireGuard VPN", "端對端加密通訊\n山區 4G LTE 穩定傳輸，延遲 ~230ms"),
]
y_s = 1.8
for title, desc in solutions:
    box(s, 7.05, y_s, 5.8, 0.8, fill=RGBColor(0x1a, 0x30, 0x25), line=C_GREEN, line_w=Pt(0.5))
    txt(s, f"◆  {title}", 7.15, y_s + 0.03, 5.6, 0.28, size=13, bold=True, color=C_GREEN)
    txt(s, desc, 7.15, y_s + 0.30, 5.6, 0.45, size=11, color=C_GRAY)
    y_s += 0.98


# ═══════════════════════════════════════════════
# Slide 3 — 系統整體架構
# ═══════════════════════════════════════════════
s = add_slide()
header_bar(s, "系統整體架構",
           "Jetson AGX Thor → WireGuard VPN → Hub Dashboard", "03")
footer(s)

# UAV 節點框
box(s, 0.2, 1.2, 7.2, 5.7, fill=C_CARD, line=C_ACCENT2, line_w=Pt(1.5))
box(s, 0.2, 1.2, 7.2, 0.38, fill=C_CARD)
txt(s, "🚁  無人機端  (Jetson AGX Thor · CUDA 13 · sm_110)",
    0.3, 1.22, 7.0, 0.35, size=12, bold=True, color=C_ACCENT2)

# 相機
box(s, 0.35, 1.72, 1.7, 0.55, fill=RGBColor(0x0f, 0x34, 0x60), line=C_ACCENT2, line_w=Pt(0.8))
txt(s, "D435_A\n640×480@15fps", 0.35, 1.72, 1.7, 0.55, size=11, color=C_WHITE, align=PP_ALIGN.CENTER)
box(s, 0.35, 2.38, 1.7, 0.55, fill=RGBColor(0x0f, 0x34, 0x60), line=C_ACCENT2, line_w=Pt(0.8))
txt(s, "D435_B\n640×480@15fps", 0.35, 2.38, 1.7, 0.55, size=11, color=C_WHITE, align=PP_ALIGN.CENTER)

# 箭頭文字
txt(s, "→", 2.1, 1.88, 0.4, 0.3, size=18, color=C_ACCENT2, align=PP_ALIGN.CENTER)
txt(s, "→", 2.1, 2.54, 0.4, 0.3, size=18, color=C_ACCENT2, align=PP_ALIGN.CENTER)

# pyrealsense2
box(s, 2.55, 1.88, 1.8, 1.0, fill=RGBColor(0x14, 0x28, 0x3e), line=C_ACCENT2, line_w=Pt(0.8))
txt(s, "pyrealsense2\nBGR Frame\nbuffer", 2.55, 1.88, 1.8, 1.0,
    size=11, color=C_WHITE, align=PP_ALIGN.CENTER)

txt(s, "→", 4.4, 2.25, 0.4, 0.3, size=18, color=C_ACCENT2, align=PP_ALIGN.CENTER)

# SAM3
box(s, 4.85, 1.72, 2.3, 1.4, fill=RGBColor(0x2c, 0x10, 0x20), line=C_ACCENT, line_w=Pt(1.2))
txt(s, "Meta SAM3\nText Grounding\nfire / smoke / person\nFP16 autocast", 4.85, 1.72, 2.3, 1.4,
    size=11, color=C_WHITE, align=PP_ALIGN.CENTER)

txt(s, "↓", 5.85, 3.2, 0.5, 0.3, size=18, color=C_ACCENT, align=PP_ALIGN.CENTER)

# MAVLink
box(s, 0.35, 3.1, 1.7, 0.9, fill=RGBColor(0x1a, 0x30, 0x15), line=C_GREEN, line_w=Pt(0.8))
txt(s, "飛控\n/dev/ttyACM1\nMAVLink 115200", 0.35, 3.1, 1.7, 0.9,
    size=10, color=C_GREEN, align=PP_ALIGN.CENTER)
txt(s, "→ GPS Cache →", 2.1, 3.35, 2.7, 0.3, size=11, color=C_GREEN, align=PP_ALIGN.CENTER)

# FSM
box(s, 4.85, 3.25, 2.3, 1.4, fill=RGBColor(0x2a, 0x25, 0x10), line=C_YELLOW, line_w=Pt(1))
txt(s, "Event State Machine\nNormal→Suspected\n→Confirmed→Tracking\n→Lost", 4.85, 3.25, 2.3, 1.4,
    size=11, color=C_YELLOW, align=PP_ALIGN.CENTER)

txt(s, "↓ should_send?", 5.85, 4.72, 1.5, 0.3, size=11, color=C_YELLOW, align=PP_ALIGN.CENTER)

# fuse_frames
box(s, 4.85, 5.1, 2.3, 0.85, fill=RGBColor(0x1a, 0x2a, 0x40), line=C_ACCENT2, line_w=Pt(0.8))
txt(s, "fuse_frames()\n雙相機拼接 + Overlay", 4.85, 5.1, 2.3, 0.85,
    size=11, color=C_WHITE, align=PP_ALIGN.CENTER)

txt(s, "→", 7.2, 5.4, 0.4, 0.3, size=18, color=C_ACCENT, align=PP_ALIGN.CENTER)

# 傳輸
box(s, 7.65, 4.8, 1.9, 1.55, fill=RGBColor(0x20, 0x10, 0x10), line=C_ACCENT, line_w=Pt(1))
txt(s, "send_alert()\nWebP thumb\nWebP evidence\nmeta JSON", 7.65, 4.8, 1.9, 1.55,
    size=10, color=C_WHITE, align=PP_ALIGN.CENTER)

# Hub 框
box(s, 9.8, 1.2, 3.3, 5.7, fill=RGBColor(0x0d, 0x1b, 0x2a), line=C_GREEN, line_w=Pt(1.5))
box(s, 9.8, 1.2, 3.3, 0.38, fill=RGBColor(0x0d, 0x1b, 0x2a))
txt(s, "🖥️  Hub 地面站  (10.0.0.7:8080)",
    9.82, 1.22, 3.2, 0.35, size=11, bold=True, color=C_GREEN)

# VPN
box(s, 9.82, 1.7, 3.1, 0.75, fill=RGBColor(0x15, 0x25, 0x15), line=C_GREEN, line_w=Pt(0.8))
txt(s, "WireGuard VPN\n10.0.0.20 → 10.0.0.7\n~230ms @ 4G LTE", 9.82, 1.7, 3.1, 0.75,
    size=10, color=C_GREEN, align=PP_ALIGN.CENTER)

box(s, 9.82, 2.65, 3.1, 0.75, fill=RGBColor(0x1a, 0x1a, 0x35), line=C_ACCENT2, line_w=Pt(0.8))
txt(s, "Flask Server\nPOST /api/v1/alerts\n驗證 Token + 存檔", 9.82, 2.65, 3.1, 0.75,
    size=10, color=C_WHITE, align=PP_ALIGN.CENTER)

box(s, 9.82, 3.58, 3.1, 0.75, fill=RGBColor(0x1a, 0x1a, 0x35), line=C_ACCENT2, line_w=Pt(0.8))
txt(s, "hub_events/<event_id>/\nmeta.json / thumb.webp\nevidence.webp", 9.82, 3.58, 3.1, 0.75,
    size=10, color=C_GRAY, align=PP_ALIGN.CENTER)

box(s, 9.82, 4.51, 3.1, 0.95, fill=RGBColor(0x1a, 0x1a, 0x35), line=C_YELLOW, line_w=Pt(0.8))
txt(s, "Web Dashboard\n自動 10s 刷新\nDETECTION / FIRE / PERSON", 9.82, 4.51, 3.1, 0.95,
    size=10, color=C_YELLOW, align=PP_ALIGN.CENTER)

box(s, 9.82, 5.64, 3.1, 0.6, fill=RGBColor(0x1a, 0x1a, 0x1a), line=C_GRAY, line_w=Pt(0.8))
txt(s, "queue/ 離線容錯備份", 9.82, 5.64, 3.1, 0.6,
    size=10, color=C_GRAY, align=PP_ALIGN.CENTER)

# VPN 箭頭
txt(s, "══════ WireGuard VPN ══════►", 7.5, 3.5, 2.7, 0.35,
    size=11, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# Slide 4 — 硬體規格
# ═══════════════════════════════════════════════
s = add_slide()
header_bar(s, "硬體規格", "計算平台 · 相機 · 飛控 · 網路", "04")
footer(s)

hw_items = [
    ("計算平台", "NVIDIA Jetson AGX Thor Developer Kit", C_ACCENT2),
    ("GPU", "NVIDIA Thor GPU (sm_110)，72-core Ampere architecture", C_ACCENT2),
    ("CUDA 版本", "CUDA 13.0 · 容器：PyTorch 2.9.0+cu130", C_ACCENT2),
    ("相機 A", "Intel RealSense D435  (SN: 332522075298)  USB3", C_ACCENT),
    ("相機 B", "Intel RealSense D435  (SN: 332522073133)  USB3", C_ACCENT),
    ("相機解析度", "640 × 480 @ 15 FPS  BGR 彩色串流 + 深度流", C_ACCENT),
    ("飛行控制器", "ArduPilot  /dev/ttyACM1  MAVLink Baud 115200", C_GREEN),
    ("GPS 訊息", "GLOBAL_POSITION_INT / GPS_RAW_INT (fix_type ≥ 2)", C_GREEN),
    ("網路", "WireGuard VPN  Thor:10.0.0.20 → Hub:10.0.0.7", C_YELLOW),
    ("VPN 延遲", "~230 ms (山區 4G LTE)  ·  峰值傳輸 ~264 Kbps", C_YELLOW),
    ("容器映像", "uav-fire-detector:latest  (24.7 GB)  pyrealsense2 2.55.1", C_GRAY),
]

cols = [hw_items[:6], hw_items[6:]]
x_starts = [0.25, 6.85]
for col_i, col_items in enumerate(cols):
    x = x_starts[col_i]
    y = 1.3
    for label, value, color in col_items:
        box(s, x, y, 6.3, 0.72, fill=C_CARD, line=color, line_w=Pt(0.8))
        txt(s, label, x + 0.12, y + 0.05, 1.8, 0.3, size=12, bold=True, color=color)
        txt(s, value, x + 0.12, y + 0.35, 6.0, 0.33, size=12, color=C_WHITE)
        y += 0.82


# ═══════════════════════════════════════════════
# Slide 5 — 軟體堆疊
# ═══════════════════════════════════════════════
s = add_slide()
header_bar(s, "軟體堆疊",
           "OS · Container · DL Framework · 通訊 · 推論", "05")
footer(s)

layers = [
    ("Layer 7  Hub 伺服器", [
        ("Flask 3.x", "REST API + Web Dashboard", C_GREEN),
        ("hub_server.py", "POST /api/v1/alerts  ·  Dashboard 自動刷新", C_GREEN),
    ], C_GREEN),
    ("Layer 6  傳輸層", [
        ("requests 2.31+", "multipart/form-data POST  ·  WebP Q60/Q75", C_ACCENT2),
        ("WireGuard", "端對端加密 VPN  ·  230ms 延遲", C_ACCENT2),
    ], C_ACCENT2),
    ("Layer 5  深度學習推論", [
        ("Meta SAM3", "Text Grounding  ·  Zero-shot  ·  FP16 autocast", C_ACCENT),
        ("PyTorch 2.9.0+cu130", "sm_110 (NVIDIA Thor)  ·  CUDA 13.0", C_ACCENT),
    ], C_ACCENT),
    ("Layer 4  相機 SDK", [
        ("pyrealsense2 2.55.1", "D435 USB3 直讀  ·  BGR 480p@15fps", C_YELLOW),
        ("pymavlink 2.4.49", "MAVLink GPS 讀取  ·  背景執行緒", C_YELLOW),
    ], C_YELLOW),
    ("Layer 3  Container Runtime", [
        ("Docker + NVIDIA CT", "uav-fire-detector:latest  24.7GB", C_GRAY),
        ("Python 3.12 venv", "/opt/venv  ·  所有相依套件獨立", C_GRAY),
    ], C_GRAY),
    ("Layer 2  OS / BSP", [
        ("Ubuntu 22.04 (Jetson BSP)", "ROS2 Jazzy  /opt/ros/jazzy", C_GRAY),
    ], C_GRAY),
]

y_l = 1.2
for layer_name, items, color in layers:
    bh = 0.38 + len(items) * 0.42
    box(s, 0.25, y_l, 12.8, bh, fill=C_CARD, line=color, line_w=Pt(1))
    box(s, 0.25, y_l, 2.5, bh, fill=RGBColor(
        int(color[0]*0.35), int(color[1]*0.35), int(color[2]*0.35)
    ))
    txt(s, layer_name, 0.35, y_l + 0.06, 2.3, bh - 0.1,
        size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
    y_i = y_l + 0.05
    for sw, desc, _ in items:
        txt(s, sw, 2.9, y_i, 2.8, 0.35, size=12, bold=True, color=color)
        txt(s, desc, 5.8, y_i, 7.0, 0.35, size=11, color=C_GRAY)
        y_i += 0.42
    y_l += bh + 0.08


# ═══════════════════════════════════════════════
# Slide 6 — SAM3 推論流程
# ═══════════════════════════════════════════════
s = add_slide()
header_bar(s, "SAM3 文字引導式偵測（Text-Grounding）",
           "Zero-shot · FP16 · 每幀 set_image 只呼叫 1 次", "06")
footer(s)

# 流程圖
steps = [
    ("D435 BGR Frame\n640 × 480", C_ACCENT2),
    ("BGR→RGB\nPIL.Image", C_ACCENT2),
    ("set_image()\n特徵提取\n(共用)", C_ACCENT),
    ("set_text_prompt()\n× 3 labels\nfire/smoke/person\n(FP16 autocast)", C_ACCENT),
    ("boxes / scores\n/ masks\n輸出", C_YELLOW),
    ("過濾\nscore≥0.45\nmask≥200px", C_GREEN),
    ("Detection List\n[label,score,\nbbox,mask]", C_GREEN),
]

x_s = 0.3
for i, (label, color) in enumerate(steps):
    bw = 1.55
    box(s, x_s, 1.5, bw, 1.6, fill=C_CARD, line=color, line_w=Pt(1.5))
    txt(s, label, x_s + 0.05, 1.6, bw - 0.1, 1.4,
        size=12, bold=(i in [2, 3]), color=color, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        txt(s, "→", x_s + bw, 2.1, 0.35, 0.5, size=20,
            color=C_ACCENT2, align=PP_ALIGN.CENTER)
    x_s += bw + 0.35

# 說明框
box(s, 0.3, 3.45, 12.5, 3.3, fill=C_CARD, line=C_ACCENT2, line_w=Pt(0.8))

details = [
    ("◆ 模型", "Meta SAM3 Image Model  (facebook/sam3)  ·  3.45 GB  ·  解析度 1008×1008  ·  信心門檻 0.45"),
    ("◆ 效能關鍵", "set_image() 每幀只呼叫 1 次（所有標籤共用特徵）  →  相較3次呼叫節省 ~66% 特徵提取時間"),
    ("◆ 精度優化", "FP16 autocast (torch.float16)  ·  只包住 PyTorch 推論區塊  ·  避免 BF16/OpenCV 型別衝突"),
    ("◆ Mask 處理", "masks tensor 形狀 (N,K,H,W) → while ndim>2: mask=mask[0] 壓縮到 2D  ·  resize 到原圖尺寸"),
    ("◆ 推論頻率", "2 Hz (每 0.5 秒)  ·  雙相機各 1 次  ·  等效每秒 4 次 set_image  ·  可調至 4Hz"),
]
y_d = 3.6
for label, desc in details:
    txt(s, label, 0.5, y_d, 1.7, 0.45, size=12, bold=True, color=C_ACCENT2)
    txt(s, desc,  2.25, y_d, 10.4, 0.45, size=11, color=C_WHITE)
    y_d += 0.5


# ═══════════════════════════════════════════════
# Slide 7 — 事件狀態機
# ═══════════════════════════════════════════════
s = add_slide()
header_bar(s, "事件狀態機（Event State Machine）",
           "5 狀態 FSM · per-label · 滑動視窗累積命中率", "07")
footer(s)

# 狀態圖
states = [
    ("NORMAL",    0.4,  3.6, C_GRAY),
    ("SUSPECTED", 2.7,  3.6, C_ORANGE),
    ("CONFIRMED", 5.0,  3.6, C_ACCENT),
    ("TRACKING",  7.3,  3.6, C_GREEN),
    ("LOST",      9.6,  3.6, C_GRAY),
]
for name, x, y, color in states:
    box(s, x, y, 2.0, 0.85, fill=C_CARD, line=color, line_w=Pt(2))
    txt(s, name, x, y + 0.18, 2.0, 0.5, size=14, bold=True,
        color=color, align=PP_ALIGN.CENTER)

# 轉移箭頭 + 條件
transitions = [
    (1.15, 3.97, "M1內≥N1\n命中"),
    (3.45, 3.97, "M2內≥N2\n命中"),
    (5.75, 3.97, "進入"),
    (8.05, 3.97, "T_LOST\n秒無命中"),
]
for x, y, label in transitions:
    txt(s, "→", x + 0.65, y + 0.05, 0.5, 0.3, size=18, color=C_ACCENT2, align=PP_ALIGN.CENTER)
    txt(s, label, x + 0.4, y + 0.33, 0.9, 0.4, size=9, color=C_ACCENT2, align=PP_ALIGN.CENTER)

# Lost → Normal 回饋
txt(s, "↩ 重置 → NORMAL", 5.0, 4.75, 2.4, 0.35, size=11, color=C_GRAY, align=PP_ALIGN.CENTER)
txt(s, "Suspected 視窗清空 ↩ NORMAL", 2.4, 4.75, 3.0, 0.35, size=10, color=C_GRAY)

# 參數表
box(s, 0.25, 1.25, 12.8, 2.1, fill=C_CARD, line=C_ACCENT2, line_w=Pt(0.8))
params = [
    ("LOOP_INTERVAL", "0.5 s", "推論間隔（2 Hz）"),
    ("SM_M1 / SM_N1", "5 幀 / 2 次", "→ Suspected：2.5s 視窗內 2 次命中"),
    ("SM_M2 / SM_N2", "10 幀 / 4 次", "→ Confirmed：5s 視窗內 4 次命中"),
    ("SM_T_SEND",     "2.0 s", "Tracking：最短傳送間隔"),
    ("SM_T_LOST",     "10.0 s", "無命中後進入 Lost 狀態"),
    ("SCORE_THR",     "0.45", "SAM3 最低信心分數門檻"),
]
cols = [params[:3], params[3:]]
for ci, col in enumerate(cols):
    xc = 0.4 + ci * 6.4
    yc = 1.4
    for name, val, desc in col:
        txt(s, name, xc, yc, 1.9, 0.35, size=11, bold=True, color=C_ACCENT2)
        txt(s, val,  xc + 1.95, yc, 0.9, 0.35, size=11, bold=True, color=C_YELLOW)
        txt(s, desc, xc + 2.9, yc, 2.9, 0.35, size=10, color=C_GRAY)
        yc += 0.42

# 時序示意
box(s, 0.25, 5.3, 12.8, 1.5, fill=C_CARD, line=C_YELLOW, line_w=Pt(0.8))
txt(s, "時序示意：", 0.4, 5.35, 1.6, 0.35, size=12, bold=True, color=C_YELLOW)
txt(s,
    "時間(s): 0    0.5   1.0   1.5   2.0   2.5   3.0   3.5   4.0   4.5   5.0   5.5   6.0\n"
    "命中:    ○     ●     ●     ○     ●     ●     ○     ●     ●     ●     ●     ●     ●\n"
    "狀態:  Normal  Susp  Susp  Susp  Susp  Conf  Track Track Track Track Track Track\n"
    "傳輸:                                         🔴                  🔴              🔴",
    0.4, 5.7, 12.4, 1.0, size=10.5, color=C_WHITE)


# ═══════════════════════════════════════════════
# Slide 8 — 頻寬分析
# ═══════════════════════════════════════════════
s = add_slide()
header_bar(s, "傳輸內容與頻寬分析",
           "事件式 vs 連續串流  ·  山區 4G LTE 可行性", "08")
footer(s)

# 傳輸內容表
box(s, 0.25, 1.25, 6.0, 3.1, fill=C_CARD, line=C_ACCENT2, line_w=Pt(1))
box(s, 0.25, 1.25, 6.0, 0.4, fill=RGBColor(0x0f, 0x34, 0x60))
txt(s, "  每次 POST 傳輸內容", 0.25, 1.25, 6.0, 0.4,
    size=13, bold=True, color=C_WHITE)
payload = [
    ("meta.json",   "事件ID / GPS / 時間戳 / bbox / label / score",  "< 1 KB",  C_ACCENT2),
    ("thumb.webp",  "雙相機拼接縮圖 320×240  WebP Q60",              "5–15 KB", C_ACCENT2),
    ("evidence.webp","ROI 截圖 + SAM3 Mask + BBox  WebP Q75",        "15–50 KB",C_ACCENT),
    ("合計",         "—",                                              "21–66 KB",C_GREEN),
]
yp = 1.75
for field, desc, size, color in payload:
    box(s, 0.35, yp, 5.8, 0.5, fill=RGBColor(0x1a, 0x24, 0x3a), line=color, line_w=Pt(0.5))
    txt(s, field, 0.45, yp + 0.07, 1.2, 0.35, size=11, bold=True, color=color)
    txt(s, desc,  1.7,  yp + 0.07, 3.0, 0.35, size=10, color=C_GRAY)
    txt(s, size,  4.75, yp + 0.07, 1.2, 0.35, size=11, bold=True, color=color)
    yp += 0.57

# 頻寬比較
box(s, 6.6, 1.25, 6.45, 3.1, fill=C_CARD, line=C_GREEN, line_w=Pt(1))
box(s, 6.6, 1.25, 6.45, 0.4, fill=RGBColor(0x0d, 0x2a, 0x1a))
txt(s, "  頻寬需求比較", 6.6, 1.25, 6.45, 0.4,
    size=13, bold=True, color=C_GREEN)
bw_items = [
    ("本系統 Tracking 峰值",  "66KB × (1/2s)",    "264 Kbps",   C_ACCENT),
    ("本系統 Tracking 均值",  "40KB × (1/3s)",    "107 Kbps",   C_YELLOW),
    ("本系統 靜默（無事件）", "0",                "0 Kbps",     C_GREEN),
    ("720p@15fps H.264 串流", "連續傳輸",          "1,000–2,000 Kbps", C_ORANGE),
    ("節省比例",               "—",               "95–99%",     C_GREEN),
]
yb = 1.75
for scenario, formula, bw, color in bw_items:
    box(s, 6.7, yb, 6.25, 0.5, fill=RGBColor(0x0d, 0x1e, 0x14), line=color, line_w=Pt(0.5))
    txt(s, scenario, 6.8, yb + 0.07, 2.8, 0.35, size=10, color=C_GRAY)
    txt(s, formula,  9.65, yb + 0.07, 1.4, 0.35, size=10, color=C_GRAY)
    txt(s, bw,       11.1, yb + 0.07, 1.5, 0.35, size=11, bold=True, color=color)
    yb += 0.57

# 可行性說明
box(s, 0.25, 4.65, 12.8, 1.55, fill=C_CARD, line=C_GREEN, line_w=Pt(1))
txt(s, "✅  山區低頻寬可行性分析", 0.4, 4.7, 4.0, 0.4,
    size=14, bold=True, color=C_GREEN)
feasibility = [
    ("4G LTE 山區（200Kbps–2Mbps）", "✅ 完全可行  ·  追蹤時均值 107 Kbps  ·  靜默時 0 Kbps"),
    ("Starlink 低延遲衛星（>5Mbps）", "✅ 完全可行  ·  可縮短 SM_T_SEND 至 1.0s 提升響應速度"),
    ("Iridium/Thuraya 衛星（64Kbps）", "⚠️ 需調整：SM_T_SEND≥8s + WebP Q30 → ~20KB/次 ≈ 20Kbps"),
]
yf = 5.15
for scenario, desc in feasibility:
    txt(s, scenario, 0.4, yf, 3.6, 0.3, size=11, bold=True, color=C_ACCENT2)
    txt(s, desc,     4.1, yf, 8.8, 0.3, size=11, color=C_WHITE)
    yf += 0.38


# ═══════════════════════════════════════════════
# Slide 9 — GPS 定位機制
# ═══════════════════════════════════════════════
s = add_slide()
header_bar(s, "GPS 定位與 MAVLink 通訊",
           "背景執行緒 · fix_type 判斷 · Home GPS 備援", "09")
footer(s)

box(s, 0.25, 1.25, 5.8, 5.6, fill=C_CARD, line=C_GREEN, line_w=Pt(1))
txt(s, "MAVLink 通訊架構", 0.4, 1.3, 5.5, 0.4, size=14, bold=True, color=C_GREEN)

mav_items = [
    ("連線方式", "/dev/ttyACM1  Baud 115200  ArduPilot"),
    ("訊息類型", "GLOBAL_POSITION_INT  (lat/lon/alt × 1e-7)"),
    ("          ", "GPS_RAW_INT           (fix_type 0–6)"),
    ("讀取架構", "獨立背景 daemon Thread  不阻塞推論主迴圈"),
    ("斷線恢復", "try/except → sleep(5) → 自動重連"),
    ("GPS 快取", "_gps_cache dict  含 fix_type + source 欄位"),
    ("精度", "單頻 GPS ±3–5m  3D Fix (fix_type=3) 有效"),
]
ym = 1.8
for label, val in mav_items:
    txt(s, label, 0.4, ym, 1.7, 0.35, size=11, bold=True, color=C_GREEN)
    txt(s, val,   2.15, ym, 3.75, 0.35, size=11, color=C_WHITE)
    ym += 0.48

box(s, 6.35, 1.25, 6.7, 5.6, fill=C_CARD, line=C_ACCENT2, line_w=Pt(1))
txt(s, "GPS fix_type 判斷邏輯", 6.5, 1.3, 6.5, 0.4,
    size=14, bold=True, color=C_ACCENT2)

fix_items = [
    ("fix_type 0", "無 GPS 定位  → 回傳 Home GPS", C_ACCENT),
    ("fix_type 1", "Dead Reckoning  → 回傳 Home GPS", C_ACCENT),
    ("fix_type 2", "2D Fix  (lat/lon)  → 使用實際 GPS ✅", C_GREEN),
    ("fix_type 3", "3D Fix  (lat/lon/alt)  → 使用實際 GPS ✅", C_GREEN),
]
yf2 = 1.85
for ft, desc, color in fix_items:
    box(s, 6.5, yf2, 6.4, 0.52, fill=RGBColor(0x0d, 0x1a, 0x2e), line=color, line_w=Pt(0.8))
    txt(s, ft, 6.62, yf2 + 0.1, 1.2, 0.3, size=12, bold=True, color=color)
    txt(s, desc, 7.9, yf2 + 0.1, 4.8, 0.3, size=11, color=C_WHITE)
    yf2 += 0.62

box(s, 6.5, 4.45, 6.4, 1.25, fill=RGBColor(0x1a, 0x20, 0x10), line=C_GREEN, line_w=Pt(0.8))
txt(s, "Home GPS 備援機制（室內測試）", 6.62, 4.5, 6.2, 0.35,
    size=12, bold=True, color=C_GREEN)
txt(s,
    "預設位置：國立清華大學光復路二段 101 號，新竹市\n"
    "座標：lat=24.7968, lon=120.9961, alt=75m\n"
    "Dashboard 自動顯示地名 + Google Maps 連結",
    6.62, 4.88, 6.2, 0.77, size=11, color=C_WHITE)


# ═══════════════════════════════════════════════
# Slide 10 — 雙相機拼接設計
# ═══════════════════════════════════════════════
s = add_slide()
header_bar(s, "雙相機拼接與視覺化設計",
           "fuse_frames() · Detection Overlay · Evidence 圖", "10")
footer(s)

# 相機示意
box(s, 0.25, 1.25, 5.9, 3.3, fill=C_CARD, line=C_ACCENT2, line_w=Pt(1))
txt(s, "雙相機拼接輸出（1284 × 480）", 0.4, 1.3, 5.6, 0.35,
    size=13, bold=True, color=C_ACCENT2)

box(s, 0.35, 1.75, 2.65, 1.85,
    fill=RGBColor(0x0d, 0x20, 0x35), line=C_GREEN, line_w=Pt(1))
txt(s, "CAM-A (d435_a)\n640×480 → 等高 resize\n前向廣角\n[person:0.98 ✅]",
    0.35, 1.75, 2.65, 1.85, size=11, color=C_GREEN, align=PP_ALIGN.CENTER)

box(s, 3.18, 1.75, 0.08, 1.85, fill=C_GRAY)  # 分隔線

box(s, 3.35, 1.75, 2.65, 1.85,
    fill=RGBColor(0x1a, 0x10, 0x10), line=C_ACCENT, line_w=Pt(1))
txt(s, "CAM-B (d435_b)\n640×480 → 等高 resize\n側向/後向\n[fire:0.94 ✅]",
    3.35, 1.75, 2.65, 1.85, size=11, color=C_ACCENT, align=PP_ALIGN.CENTER)

overlay_items = [
    ("mask overlay",  "目標區域半透明著色  green=person / red=fire"),
    ("bbox rectangle","偵測框 + label + score 文字"),
    ("相機標籤",       "CAM-A(d435_a) / CAM-B(d435_b) 左上角"),
    ("分隔線",         "4px 灰色垂直分隔線"),
]
yo = 3.75
for label, desc in overlay_items:
    txt(s, f"◆ {label}", 0.4, yo, 2.0, 0.35, size=11, bold=True, color=C_ACCENT2)
    txt(s, desc,         2.5, yo, 3.5, 0.35, size=11, color=C_GRAY)
    yo += 0.42

# 右側
box(s, 6.5, 1.25, 6.6, 5.6, fill=C_CARD, line=C_YELLOW, line_w=Pt(1))
txt(s, "傳輸圖像說明", 6.65, 1.3, 6.4, 0.35, size=13, bold=True, color=C_YELLOW)

img_types = [
    ("thumb.webp", "WebP Q60  ·  320×240 px  ·  ~9 KB",
     "全景拼接縮圖  ·  快速預覽用\n包含雙相機視角與所有偵測框", C_ACCENT2),
    ("evidence.webp", "WebP Q75  ·  原始尺寸  ·  ~48 KB",
     "雙相機全景 + SAM3 Mask overlay\n所有偵測框標籤 (bbox 覆蓋>85% 跳過重繪)\n為 Dashboard 主要顯示圖像", C_YELLOW),
]
yi = 1.75
for fname, spec, desc, color in img_types:
    box(s, 6.6, yi, 6.4, 2.3, fill=RGBColor(0x1a, 0x1e, 0x10), line=color, line_w=Pt(1))
    txt(s, fname, 6.75, yi + 0.1, 4.5, 0.38, size=14, bold=True, color=color)
    txt(s, spec,  6.75, yi + 0.5, 6.1, 0.35, size=11, color=C_ACCENT2)
    txt(s, desc,  6.75, yi + 0.85, 6.1, 1.0, size=11, color=C_GRAY)
    yi += 2.5

# Note
box(s, 0.25, 4.85, 5.9, 1.5, fill=RGBColor(0x1a, 0x1a, 0x10), line=C_YELLOW, line_w=Pt(0.8))
txt(s, "Note 欄位格式", 0.4, 4.9, 3.0, 0.35, size=12, bold=True, color=C_YELLOW)
txt(s, "camA=[person:0.99, person:0.98] camB=[fire:0.94, person:0.92] sm=tracking",
    0.4, 5.28, 5.6, 0.35, size=10.5, color=C_WHITE)
txt(s, "各相機偵測結果摘要（依信心分數降排）+ 當前狀態機狀態",
    0.4, 5.65, 5.6, 0.35, size=10.5, color=C_GRAY)


# ═══════════════════════════════════════════════
# Slide 11 — Hub Dashboard
# ═══════════════════════════════════════════════
s = add_slide()
header_bar(s, "Hub 地面站 Dashboard",
           "Flask REST API · 自動刷新 · 地名解析 · 點圖開原圖", "11")
footer(s)

box(s, 0.25, 1.25, 4.3, 5.6, fill=C_CARD, line=C_ACCENT2, line_w=Pt(1))
txt(s, "REST API 端點", 0.4, 1.3, 4.0, 0.38, size=13, bold=True, color=C_ACCENT2)

api_items = [
    ("POST", "/api/v1/alerts", "接收告警：Token 驗證 → 存 meta.json\n+ thumb.webp + evidence.webp"),
    ("GET",  "/api/v1/pull_requests", "Reachback 查詢（保留擴充用）"),
    ("GET",  "/",               "Web Dashboard 自動 10s 刷新"),
    ("GET",  "/events/<id>/<f>","靜態圖片直接讀取服務"),
]
ya = 1.8
for method, path, desc in api_items:
    mc = C_GREEN if method == "GET" else C_ACCENT
    box(s, 0.35, ya, 3.9, 0.88, fill=RGBColor(0x0d, 0x1e, 0x2e), line=mc, line_w=Pt(0.5))
    box(s, 0.35, ya, 0.55, 0.88, fill=mc)
    txt(s, method, 0.35, ya + 0.25, 0.55, 0.38, size=10, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, path, 0.95, ya + 0.07, 3.2, 0.35, size=10, bold=True, color=mc)
    txt(s, desc, 0.95, ya + 0.43, 3.2, 0.42, size=9.5, color=C_GRAY)
    ya += 1.02

box(s, 4.8, 1.25, 8.3, 5.6, fill=C_CARD, line=C_YELLOW, line_w=Pt(1))
txt(s, "Dashboard 功能特色", 4.95, 1.3, 8.0, 0.38, size=13, bold=True, color=C_YELLOW)

dash_features = [
    ("DETECTION / FIRE 標題",
     "從 Note 欄解析最高優先標籤\nfire > smoke > person  ·  顏色編碼 red/gray/green", C_ACCENT),
    ("GPS 地名解析",
     "座標對照 → 自動顯示中文地名\n國立清華大學/合歡山等  +  Google Maps 超連結", C_GREEN),
    ("全幅偵測圖",
     "evidence.webp 全寬顯示（最大 960px）\n點圖開新分頁顯示原始解析度", C_YELLOW),
    ("事件排序",
     "最新事件在最上方  ·  最多顯示 50 筆\nhub_events/ 目錄長期保存所有事件", C_ACCENT2),
    ("自動刷新",
     "<meta http-equiv=refresh content=10>\n無需手動重載即可看到最新事件", C_GRAY),
]
yd = 1.8
for title, desc, color in dash_features:
    box(s, 4.95, yd, 8.0, 0.92, fill=RGBColor(0x14, 0x1a, 0x10), line=color, line_w=Pt(0.5))
    txt(s, f"◆ {title}", 5.05, yd + 0.06, 7.8, 0.3, size=12, bold=True, color=color)
    txt(s, desc, 5.05, yd + 0.38, 7.8, 0.5, size=10.5, color=C_GRAY)
    yd += 1.03


# ═══════════════════════════════════════════════
# Slide 12 — 離線容錯 + 安全機制
# ═══════════════════════════════════════════════
s = add_slide()
header_bar(s, "離線容錯機制與系統安全",
           "queue/ 本機備份 · Token 認證 · MAVLink 斷線重連", "12")
footer(s)

# 左欄
box(s, 0.25, 1.25, 6.0, 5.6, fill=C_CARD, line=C_ORANGE, line_w=Pt(1))
txt(s, "離線容錯機制", 0.4, 1.3, 5.7, 0.38, size=14, bold=True, color=C_ORANGE)

offline_items = [
    ("觸發條件",
     "POST 失敗（VPN 斷線 / Hub 離線）或\nHTTP status ≠ 200 時自動觸發"),
    ("儲存內容",
     "queue/<event_id>/meta.json\nqueue/<event_id>/thumb.webp\nqueue/<event_id>/evidence.webp"),
    ("重送機制",
     "目前為手動重送（計劃實作自動重試）\n網路恢復後可 rsync 或 POST replay"),
    ("本地監控",
     "queue/ 目錄可用 ls -la 確認積壓量\n每筆約 21–66 KB，32GB SD 可存約 50 萬筆"),
]
yo2 = 1.8
for label, desc in offline_items:
    box(s, 0.35, yo2, 5.8, 1.15, fill=RGBColor(0x2a, 0x18, 0x05), line=C_ORANGE, line_w=Pt(0.5))
    txt(s, f"◆ {label}", 0.45, yo2 + 0.07, 5.6, 0.32, size=12, bold=True, color=C_ORANGE)
    txt(s, desc, 0.45, yo2 + 0.42, 5.6, 0.68, size=11, color=C_WHITE)
    yo2 += 1.28

# 右欄
box(s, 6.6, 1.25, 6.45, 5.6, fill=C_CARD, line=C_YELLOW, line_w=Pt(1))
txt(s, "系統安全機制", 6.75, 1.3, 6.2, 0.38, size=14, bold=True, color=C_YELLOW)

security_items = [
    ("API Token 認證",
     "每次 POST 需帶 X-Auth-Token Header\n不符則回 401 Unauthorized"),
    ("WireGuard 加密",
     "Curve25519 / ChaCha20-Poly1305\nVPN 隧道端對端加密"),
    ("Docker 隔離",
     "--runtime=nvidia --privileged\n容器內獨立 venv 不影響宿主系統"),
    ("MAVLink 斷線重連",
     "背景執行緒 try/except → sleep(5)\n自動重新建立 MAVLink 連線"),
    ("GPS 備援",
     "fix_type < 2 → Home GPS\n保證告警 meta 中永遠有座標資訊"),
]
ys = 1.8
for label, desc in security_items:
    box(s, 6.7, ys, 6.25, 0.93, fill=RGBColor(0x1e, 0x1a, 0x06), line=C_YELLOW, line_w=Pt(0.5))
    txt(s, f"◆ {label}", 6.8, ys + 0.06, 6.0, 0.3, size=12, bold=True, color=C_YELLOW)
    txt(s, desc, 6.8, ys + 0.4, 6.0, 0.5, size=11, color=C_WHITE)
    ys += 1.05


# ═══════════════════════════════════════════════
# Slide 13 — 系統驗證結果
# ═══════════════════════════════════════════════
s = add_slide()
header_bar(s, "系統驗證結果（2026-02-25 室內測試）",
           "SAM3 ✅  MAVLink ✅  雙 D435 ✅  Hub Dashboard ✅", "13")
footer(s)

# 啟動序列
box(s, 0.25, 1.25, 12.8, 2.85, fill=C_CARD, line=C_ACCENT2, line_w=Pt(0.8))
txt(s, "系統啟動序列 Log", 0.4, 1.3, 6.0, 0.38, size=13, bold=True, color=C_ACCENT2)

logs = [
    ("[SAM3] ✅ 模型載入完成  3.45 GB  CUDA fp16  device=cuda", C_GREEN),
    ("[MAV]  ✅ Heartbeat OK  sysid=1  /dev/ttyACM1 baud=115200", C_GREEN),
    ("[RS]   ✅ 偵測到 2 顆 RealSense: ['332522075298', '332522073133']", C_GREEN),
    ("[RS]   ✅ 相機開啟完成，開始推論迴圈（2 Hz）", C_GREEN),
    ("[EVENT] 🚨 label=fire  level=tracking  event_id=fire_1771995086333  GPS=(24.7968,120.9961,75.0m)", C_ACCENT),
    ("[send_alert] ✅ 成功！status=200  耗時=0.27s  thumb=9.7KB  evid=48.8KB", C_YELLOW),
]
yl = 1.78
for log_text, color in logs:
    box(s, 0.35, yl, 12.5, 0.34, fill=RGBColor(0x08, 0x12, 0x1e))
    txt(s, log_text, 0.45, yl + 0.03, 12.3, 0.28, size=10.5, color=color)
    yl += 0.36

# 效能量測
box(s, 0.25, 4.35, 6.0, 2.4, fill=C_CARD, line=C_YELLOW, line_w=Pt(1))
txt(s, "實測效能數據", 0.4, 4.4, 5.7, 0.38, size=13, bold=True, color=C_YELLOW)
perf = [
    ("SAM3 模型載入",  "~52 秒（首次下載 3.45GB + 初始化）"),
    ("POST 耗時",      "0.22–0.33 秒（含 WireGuard 230ms）"),
    ("thumb 大小",     "9.4–9.7 KB（320×240 WebP Q60）"),
    ("evidence 大小",  "47.4–49.5 KB（全景 WebP Q75）"),
    ("信心分數",       "fire 0.94–0.97 / person 0.92–0.99"),
]
yp2 = 4.85
for label, val in perf:
    txt(s, label, 0.4, yp2, 2.3, 0.3, size=11, bold=True, color=C_ACCENT2)
    txt(s, val,   2.75, yp2, 3.35, 0.3, size=11, color=C_WHITE)
    yp2 += 0.36

# 驗證截圖說明
box(s, 6.6, 4.35, 6.45, 2.4, fill=C_CARD, line=C_GREEN, line_w=Pt(1))
txt(s, "Dashboard 驗證（截圖）", 6.75, 4.4, 6.2, 0.38, size=13, bold=True, color=C_GREEN)
dash_val = [
    "DETECTION FIRE  /  DETECTION PERSON",
    "Confidence: 0.9736 / 0.9858 / 0.9893",
    "GPS: lat=24.7968 lon=120.9961 alt=75.0m",
    "地名：國立清華大學，新竹市東區光復路二段",
    "Note: camA=[person:0.99] camB=[fire:0.94, person:0.92]",
    "sm=tracking  ·  Camera: dual  ·  ROI: [0,0,1284,480]",
]
yv = 4.85
for item in dash_val:
    txt(s, item, 6.75, yv, 6.2, 0.3, size=11, color=C_WHITE)
    yv += 0.35


# ═══════════════════════════════════════════════
# Slide 14 — 優化方向與未來工作
# ═══════════════════════════════════════════════
s = add_slide()
header_bar(s, "優化方向與未來工作",
           "推論加速 · 傳輸優化 · 偵測精度 · 定位增強", "14")
footer(s)

opt_cols = [
    ("🚀 推論速度優化", C_ACCENT, [
        ("TensorRT INT8 量化", "torch2trt  ·  預期 3–5× 加速（Thor GPU 原生支援）"),
        ("縮小輸入解析度", "1008→512 px  ·  推論時間 -60%  ·  精度略降"),
        ("非同步雙相機推論", "A/B 各開獨立執行緒  ·  延遲降低 50%"),
        ("提升推論頻率", "0.5s→0.25s (4Hz)  ·  響應更快"),
    ]),
    ("📡 傳輸頻寬優化", C_ACCENT2, [
        ("動態 WebP 品質", "依頻寬自動調 Q10–Q75  ·  衛星環境壓到 5KB/次"),
        ("差異幀傳輸", "只傳與前次事件的差異區域  ·  減少 40–70%"),
        ("事件等級分流", "fire=高頻寬 / person=低頻寬  ·  資源優先"),
        ("SM_T_SEND 動態調整", "依網路品質動態調整傳送間隔"),
    ]),
    ("🎯 偵測精度優化", C_GREEN, [
        ("IoU 跨幀追蹤", "SORT/ByteTrack  ·  更穩定 Tracking 狀態"),
        ("D435 深度資訊", "啟用深度流過濾近距離誤偵測目標"),
        ("分類別信心門檻", "fire=0.5 / smoke=0.6 / person=0.4 獨立調整"),
        ("False Positive 分析", "收集室外飛行誤報資料，調整 SM 參數"),
    ]),
    ("📍 定位精度優化", C_YELLOW, [
        ("RTK GPS 整合", "搭配地面站 RTK 基站  ·  精度 < 10 cm"),
        ("氣壓計高度融合", "無 GPS Fix 時用 BARO_ALT 補充高度"),
        ("時間戳精確對齊", "MAVLink 時間戳與系統時鐘 NTP 同步"),
        ("geofence 告警", "設定禁飛區邊界自動告警"),
    ]),
]

x_start = 0.2
for col_i, (title, color, items) in enumerate(opt_cols):
    xc = x_start + col_i * 3.25
    box(s, xc, 1.2, 3.1, 5.7, fill=C_CARD, line=color, line_w=Pt(1))
    box(s, xc, 1.2, 3.1, 0.42, fill=RGBColor(
        int(color[0]*0.5), int(color[1]*0.5), int(color[2]*0.5)
    ))
    txt(s, title, xc + 0.05, 1.22, 3.0, 0.4, size=11.5, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER)
    yi = 1.75
    for item_title, item_desc in items:
        box(s, xc + 0.1, yi, 2.9, 1.12, fill=RGBColor(0x10, 0x18, 0x28), line=color, line_w=Pt(0.5))
        txt(s, item_title, xc + 0.18, yi + 0.06, 2.72, 0.35, size=11, bold=True, color=color)
        txt(s, item_desc,  xc + 0.18, yi + 0.42, 2.72, 0.64, size=9.5, color=C_GRAY)
        yi += 1.2


# ═══════════════════════════════════════════════
# Slide 15 — 程式碼結構
# ═══════════════════════════════════════════════
s = add_slide()
header_bar(s, "程式碼結構",
           "uav_hub/ · 主程式 ~1000 行 · 模組化設計", "15")
footer(s)

box(s, 0.25, 1.25, 5.6, 5.6, fill=C_CARD, line=C_ACCENT2, line_w=Pt(1))
txt(s, "檔案結構", 0.4, 1.3, 5.3, 0.38, size=13, bold=True, color=C_ACCENT2)

files = [
    ("thor_dualcam_event_sender.py", "主程式  ~1000 行", C_ACCENT2,
     "init_sam3() · _infer_sam3_real() · EventStateMachine\n_mav_reader_thread() · run_realsense_mode() · _main_loop_body()"),
    ("thor_send_alert.py", "傳輸函式庫", C_ACCENT2,
     "send_alert() · make_thumb() · make_evidence()\nbuild_meta() · save_to_queue()"),
    ("hub_server.py", "Hub Flask 伺服器", C_GREEN,
     "POST /api/v1/alerts · GET / Dashboard\n_location_label() · _primary_display_label()"),
    ("run_uav_alert.sh", "Docker 一鍵啟動", C_YELLOW,
     "GPU 確認 · SERIAL_DEVICES 掃描\npymavlink auto-install · 參數透傳"),
    ("decord.py", "SAM3 訓練依賴 stub", C_GRAY,
     "VideoReader/cpu() stub 繞過 ImportError\n推論時不需要真實 decord"),
    ("requirements.txt", "依賴清單", C_GRAY,
     "requests / Pillow / opencv-python / numpy"),
]
yf = 1.8
for fname, role, color, desc in files:
    box(s, 0.35, yf, 5.4, 0.8, fill=RGBColor(0x0d, 0x18, 0x2e), line=color, line_w=Pt(0.5))
    txt(s, fname, 0.45, yf + 0.05, 3.5, 0.3, size=11, bold=True, color=color)
    txt(s, role,  4.0,  yf + 0.05, 1.65, 0.3, size=10, color=C_ACCENT2)
    txt(s, desc,  0.45, yf + 0.38, 5.2, 0.38, size=9.5, color=C_GRAY)
    yf += 0.9

# 右欄 - 主程式模組說明
box(s, 6.1, 1.25, 7.0, 5.6, fill=C_CARD, line=C_ACCENT, line_w=Pt(1))
txt(s, "thor_dualcam_event_sender.py 模組說明", 6.25, 1.3, 6.7, 0.38,
    size=13, bold=True, color=C_ACCENT)

modules = [
    ("init_sam3()", "SAM3 模型初始化\nHF Token 登入 · bpe 路徑自動偵測 · CUDA 確認", C_ACCENT),
    ("_infer_sam3_real()", "實際 FP16 推論\nset_image×1 → set_text_prompt×3 → 結果解析", C_ACCENT),
    ("EventStateMachine", "per-label 5狀態 FSM\n滑動視窗命中率計算 · 節流傳送控制", C_YELLOW),
    ("_mav_reader_thread()", "MAVLink 背景執行緒\nGPS_RAW_INT/GLOBAL_POSITION_INT 快取", C_GREEN),
    ("run_realsense_mode()", "pyrealsense2 雙相機主迴圈\n序號自動偵測 · 3s 超時容錯", C_ACCENT2),
    ("_main_loop_body()", "每 0.5s 核心邏輯\n推論→FSM→拼接→GPS→POST→queue", C_ACCENT2),
    ("fuse_frames()", "雙相機拼接 + Detection overlay\nmask/bbox/label 繪製 · 480p 等高縮放", C_ACCENT2),
]
ym2 = 1.8
for func, desc, color in modules:
    box(s, 6.2, ym2, 6.8, 0.69, fill=RGBColor(0x10, 0x14, 0x22), line=color, line_w=Pt(0.5))
    txt(s, func, 6.3, ym2 + 0.05, 2.5, 0.28, size=11, bold=True, color=color)
    txt(s, desc, 6.3, ym2 + 0.35, 6.6, 0.3, size=9.5, color=C_GRAY)
    ym2 += 0.76


# ═══════════════════════════════════════════════
# Slide 16 — 總結
# ═══════════════════════════════════════════════
s = add_slide()
header_bar(s, "總結與貢獻", "研究成果 · 論文貢獻點 · 系統驗證", "16")
footer(s)

box(s, 0.25, 1.25, 8.0, 5.6, fill=C_CARD, line=C_ACCENT2, line_w=Pt(1))
txt(s, "系統貢獻摘要", 0.4, 1.3, 7.7, 0.38, size=14, bold=True, color=C_ACCENT2)

contributions = [
    ("1", "事件式傳輸架構",
     "FSM 5狀態機過濾假陽性\n相較連續串流節省 95–99% 頻寬\n山區 4G LTE 環境完全可行", C_ACCENT),
    ("2", "SAM3 Zero-shot 部署",
     "Meta SAM3 Text-Grounding\n無需針對火災/人員重新訓練\nFP16 加速 + CUDA 13 sm_110 相容", C_YELLOW),
    ("3", "雙目視覺 + MAVLink 整合",
     "雙 D435 同步採集 + 拼接視覺化\nMAVLink GPS 時間戳精確定位\nfix_type 判斷 + Home GPS 備援", C_GREEN),
    ("4", "完整部署驗證",
     "Docker 容器化一鍵部署\n室內完整驗證（D435×2 + 飛控）✅\nHub Dashboard 實際截圖證明", C_ACCENT2),
]
yc = 1.8
for num, title, desc, color in contributions:
    box(s, 0.35, yc, 7.8, 1.18, fill=RGBColor(0x0d, 0x18, 0x2e), line=color, line_w=Pt(1))
    box(s, 0.35, yc, 0.48, 1.18, fill=color)
    txt(s, num, 0.35, yc + 0.35, 0.48, 0.48, size=22, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, 0.9, yc + 0.08, 6.8, 0.38, size=13, bold=True, color=color)
    txt(s, desc, 0.9, yc + 0.48, 7.1, 0.65, size=11, color=C_GRAY)
    yc += 1.3

# 右側 Q&A
box(s, 8.55, 1.25, 4.6, 5.6, fill=C_CARD, line=C_ACCENT, line_w=Pt(1))
box(s, 8.55, 1.25, 4.6, 0.45, fill=C_ACCENT)
txt(s, "系統規格一覽", 8.7, 1.27, 4.4, 0.4, size=13, bold=True, color=C_WHITE)

specs = [
    ("推論模型",    "Meta SAM3 (3.45 GB)"),
    ("推論頻率",    "2 Hz (0.5s/次)"),
    ("偵測類別",    "fire / smoke / person"),
    ("信心門檻",    "score ≥ 0.45"),
    ("首次告警",    "~3.0 秒後確認"),
    ("最高傳送頻率","1 次 / 2 秒"),
    ("每次傳輸",    "21–66 KB"),
    ("峰值頻寬",    "264 Kbps"),
    ("靜默頻寬",    "0 Kbps"),
    ("VPN 延遲",   "~230 ms"),
    ("相機",       "D435 × 2  640×480@15fps"),
    ("GPU 平台",   "NVIDIA Thor sm_110"),
]
ys2 = 1.8
for label, val in specs:
    txt(s, label, 8.65, ys2, 1.8, 0.32, size=10.5, bold=True, color=C_ACCENT2)
    txt(s, val,   10.5, ys2, 2.55, 0.32, size=10.5, color=C_WHITE)
    ys2 += 0.37

# 儲存
out_path = "/home/alan/xin/uav_hub/UAV_Alert_System_Presentation.pptx"
prs.save(out_path)
print(f"✅ PPT 已儲存：{out_path}")
print(f"   共 {len(prs.slides)} 張投影片")
